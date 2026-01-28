"""
Orion Ingestion Module

Main interface for all document ingestion operations:
- Document loading and preprocessing
- Text chunking strategies
- Embedding generation
- Vector store management
- Batch processing with progress tracking
- Support for 30+ file formats (PDF, DOCX, TXT, MD, code files, etc.)

The ingestion pipeline orchestrates the entire process from raw documents
to searchable vector embeddings stored in ChromaDB.
"""

import hashlib
import json
import re
from concurrent.futures import ThreadPoolExecutor, as_completed
from dataclasses import dataclass, field
from datetime import datetime
from difflib import SequenceMatcher
from pathlib import Path
from typing import TYPE_CHECKING, Any, Optional

import fitz  # PyMuPDF
from langchain.schema import Document
from langchain_community.document_loaders import TextLoader
from langchain_text_splitters import RecursiveCharacterTextSplitter
from tqdm import tqdm

# Local imports
from src.retrieval.embeddings import EmbeddingManager
from src.retrieval.vector_store import ChromaVectorStore, FileType, create_vector_store
from src.utilities.config import ChunkerType
from src.utilities.utils import ensure_config, log_debug, log_error, log_info, log_warning

if TYPE_CHECKING:
    from src.utilities.config import OrionConfig


# ========== CUSTOM DOCUMENT LOADERS ==========
class PyMuPDFLoader:
    """Custom PyMuPDF loader for better PDF text extraction"""

    def __init__(self, file_path: str):
        self.file_path = file_path

    def load(self) -> list[Document]:
        """Load PDF using PyMuPDF with better text extraction"""
        documents = []

        try:
            pdf_document = fitz.open(self.file_path)

            for page_num in range(len(pdf_document)):
                page = pdf_document.load_page(page_num)

                # Extract text with better spacing preservation
                text = page.get_text("text")

                # If text is empty or very short, try alternative extraction
                if not text or len(text.strip()) < 50:
                    blocks = page.get_text("blocks")
                    text_blocks = []
                    for block in blocks:
                        if len(block) >= 5 and block[4]:
                            text_blocks.append(block[4])
                    text = "\n".join(text_blocks)

                if text.strip():
                    doc = Document(
                        page_content=text,
                        metadata={
                            "page": page_num + 1,
                            "source": self.file_path,
                            "extraction_method": "pymupdf",
                        },
                    )
                    documents.append(doc)

            pdf_document.close()

        except Exception as e:
            raise Exception(f"Failed to extract text with PyMuPDF: {e}")

        return documents


class DOCXLoader:
    """Custom DOCX loader using python-docx"""

    def __init__(self, file_path: str):
        self.file_path = file_path

    def load(self) -> list[Document]:
        """Load DOCX file"""
        try:
            import docx

            doc = docx.Document(self.file_path)
            text = "\n".join([paragraph.text for paragraph in doc.paragraphs if paragraph.text.strip()])

            return [
                Document(
                    page_content=text,
                    metadata={"source": self.file_path, "extraction_method": "python-docx"},
                )
            ]
        except ImportError:
            log_warning(
                "python-docx not installed. Install with: pip install python-docx",
                config=None,
            )
            raise Exception("python-docx not available for DOCX loading")
        except Exception as e:
            raise Exception(f"Failed to load DOCX file: {e}")


class JSONLoader:
    """Custom JSON loader"""

    def __init__(self, file_path: str):
        self.file_path = file_path

    def load(self) -> list[Document]:
        """Load JSON file"""
        try:
            with open(self.file_path, "r", encoding="utf-8") as f:
                data = json.load(f)

            # Convert JSON to formatted text
            if isinstance(data, dict):
                text = json.dumps(data, indent=2)
            elif isinstance(data, list):
                text = "\n\n".join([json.dumps(item, indent=2) for item in data])
            else:
                text = str(data)

            return [
                Document(
                    page_content=text,
                    metadata={"source": self.file_path, "extraction_method": "json"},
                )
            ]
        except Exception as e:
            raise Exception(f"Failed to load JSON file: {e}")


class MarkdownLoader:
    """Custom Markdown loader"""

    def __init__(self, file_path: str):
        self.file_path = file_path

    def load(self) -> list[Document]:
        """Load Markdown file"""
        try:
            with open(self.file_path, "r", encoding="utf-8") as f:
                text = f.read()

            return [
                Document(
                    page_content=text,
                    metadata={"source": self.file_path, "extraction_method": "markdown"},
                )
            ]
        except Exception as e:
            raise Exception(f"Failed to load Markdown file: {e}")


class CSVLoader:
    """Custom CSV loader"""

    def __init__(self, file_path: str):
        self.file_path = file_path

    def load(self) -> list[Document]:
        """Load CSV file"""
        try:
            import csv

            with open(self.file_path, "r", encoding="utf-8") as f:
                reader = csv.reader(f)
                rows = list(reader)

            # Format as text
            text = "\n".join([", ".join(row) for row in rows])

            return [
                Document(
                    page_content=text,
                    metadata={"source": self.file_path, "extraction_method": "csv"},
                )
            ]
        except Exception as e:
            raise Exception(f"Failed to load CSV file: {e}")


class PPTXLoader:
    """Custom PowerPoint loader using python-pptx"""
    
    def __init__(self, file_path: str):
        self.file_path = file_path
    
    def load(self) -> list[Document]:
        """Load PPTX file and extract text from slides"""
        try:
            from pptx import Presentation
            
            prs = Presentation(self.file_path)
            all_text = []
            
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        slide_text.append(shape.text)
                
                if slide_text:
                    all_text.append(f"=== Slide {slide_num} ===\n" + "\n".join(slide_text))
            
            text = "\n\n".join(all_text)
            
            return [
                Document(
                    page_content=text,
                    metadata={
                        "source": self.file_path,
                        "extraction_method": "python-pptx",
                        "slide_count": len(prs.slides)
                    },
                )
            ]
        except ImportError:
            log_warning("python-pptx not installed. Install with: pip install python-pptx", config=None)
            raise Exception("python-pptx not available for PPTX loading")
        except Exception as e:
            raise Exception(f"Failed to load PPTX file: {e}")


class XLSXLoader:
    """Custom Excel loader using openpyxl"""
    
    def __init__(self, file_path: str):
        self.file_path = file_path
    
    def load(self) -> list[Document]:
        """Load XLSX file and extract data from sheets"""
        try:
            from openpyxl import load_workbook
            
            wb = load_workbook(self.file_path, data_only=True)
            all_text = []
            
            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                sheet_text = [f"=== Sheet: {sheet_name} ==="]
                
                for row in sheet.iter_rows(values_only=True):
                    # Filter out None values and convert to strings
                    row_data = [str(cell) for cell in row if cell is not None]
                    if row_data:
                        sheet_text.append(", ".join(row_data))
                
                if len(sheet_text) > 1:  # More than just the header
                    all_text.append("\n".join(sheet_text))
            
            text = "\n\n".join(all_text)
            
            return [
                Document(
                    page_content=text,
                    metadata={
                        "source": self.file_path,
                        "extraction_method": "openpyxl",
                        "sheet_count": len(wb.sheetnames)
                    },
                )
            ]
        except ImportError:
            log_warning("openpyxl not installed. Install with: pip install openpyxl", config=None)
            raise Exception("openpyxl not available for XLSX loading")
        except Exception as e:
            raise Exception(f"Failed to load XLSX file: {e}")


class RTFLoader:
    """Custom RTF loader using striprtf"""
    
    def __init__(self, file_path: str):
        self.file_path = file_path
    
    def load(self) -> list[Document]:
        """Load RTF file"""
        try:
            from striprtf.striprtf import rtf_to_text
            
            with open(self.file_path, "r", encoding="utf-8", errors="ignore") as f:
                rtf_content = f.read()
            
            text = rtf_to_text(rtf_content)
            
            return [
                Document(
                    page_content=text,
                    metadata={"source": self.file_path, "extraction_method": "striprtf"},
                )
            ]
        except ImportError:
            log_warning("striprtf not installed. Install with: pip install striprtf", config=None)
            raise Exception("striprtf not available for RTF loading")
        except Exception as e:
            raise Exception(f"Failed to load RTF file: {e}")


# ========== SUPPORTED FILE TYPES ==========
# Expanded to match FileType enum from vector_store.py
SUPPORTED_EXTENSIONS = {
    # Documents
    ".pdf": "PDF Document",
    ".docx": "Word Document",
    ".doc": "Word Document (Legacy)",
    ".pptx": "PowerPoint Presentation",
    ".xlsx": "Excel Spreadsheet",
    ".txt": "Text File",
    ".md": "Markdown File",
    ".csv": "CSV File",
    ".json": "JSON File",
    ".xml": "XML File",
    ".yaml": "YAML File",
    ".yml": "YAML File",
    ".rtf": "Rich Text Format",
    # Code files
    ".py": "Python Source",
    ".js": "JavaScript Source",
    ".ts": "TypeScript Source",
    ".java": "Java Source",
    ".cpp": "C++ Source",
    ".c": "C Source",
    ".h": "C Header",
    ".hpp": "C++ Header",
    ".cs": "C# Source",
    ".go": "Go Source",
    ".rs": "Rust Source",
    ".rb": "Ruby Source",
    ".php": "PHP Source",
    ".swift": "Swift Source",
    ".kt": "Kotlin Source",
    # Web files
    ".html": "HTML File",
    ".css": "CSS File",
    ".scss": "SCSS File",
    ".jsx": "JSX File",
    ".tsx": "TSX File",
    ".vue": "Vue Component",
    # Config files
    ".ini": "INI Config",
    ".conf": "Config File",
    ".toml": "TOML Config",
}

DOCUMENT_LOADERS = {
    ".pdf": PyMuPDFLoader,
    ".docx": DOCXLoader,
    ".pptx": PPTXLoader,
    ".xlsx": XLSXLoader,
    ".rtf": RTFLoader,
    ".txt": TextLoader,
    ".md": MarkdownLoader,
    ".json": JSONLoader,
    ".csv": CSVLoader,
    # Code files use TextLoader
    ".py": TextLoader,
    ".js": TextLoader,
    ".ts": TextLoader,
    ".java": TextLoader,
    ".cpp": TextLoader,
    ".c": TextLoader,
    ".html": TextLoader,
    ".css": TextLoader,
    ".yaml": TextLoader,
    ".yml": TextLoader,
    ".xml": TextLoader,
    ".ini": TextLoader,
    ".toml": TextLoader,
    ".jsx": TextLoader,
    ".tsx": TextLoader,
    ".vue": TextLoader,
    ".go": TextLoader,
    ".rs": TextLoader,
    ".rb": TextLoader,
    ".php": TextLoader,
    ".swift": TextLoader,
    ".kt": TextLoader,
}


# ========== DATA STRUCTURES ==========
@dataclass
class IngestionStats:
    """Statistics from ingestion process"""

    total_files: int = 0
    successful_files: int = 0
    failed_files: int = 0
    total_documents: int = 0
    total_chunks: int = 0
    total_embeddings: int = 0
    processing_time: float = 0.0
    errors: list[str] = field(default_factory=list)

    @property
    def success_rate(self) -> float:
        """Calculate success rate"""
        if self.total_files == 0:
            return 0.0
        return (self.successful_files / self.total_files) * 100

    def add_error(self, error: str) -> None:
        """Add an error message"""
        if self.errors is None:
            self.errors = []
        self.errors.append(error)

    def summary(self) -> dict[str, Any]:
        """Get summary dictionary"""
        return {
            "total_files": self.total_files,
            "successful_files": self.successful_files,
            "failed_files": self.failed_files,
            "success_rate": f"{self.success_rate:.1f}%",
            "total_documents": self.total_documents,
            "total_chunks": self.total_chunks,
            "total_embeddings": self.total_embeddings,
            "processing_time": f"{self.processing_time:.2f}s",
            "errors": self.errors,
        }


@dataclass
class DocumentMetadata:
    """Metadata for ingested documents"""

    file_path: str
    file_name: str
    file_type: str
    file_size: int
    file_hash: str
    chunk_count: int
    ingestion_timestamp: str
    source_directory: str = ""

    def to_dict(self) -> dict[str, Any]:
        """Convert to dictionary"""
        return {
            "file_path": self.file_path,
            "file_name": self.file_name,
            "file_type": self.file_type,
            "file_size": self.file_size,
            "file_hash": self.file_hash,
            "chunk_count": self.chunk_count,
            "ingestion_timestamp": self.ingestion_timestamp,
            "source_directory": self.source_directory,
        }


# ========== DOCUMENT PROCESSING ==========
class DocumentProcessor:
    """Handles document loading and preprocessing"""

    def __init__(self, config: Optional["OrionConfig"] = None):
        self.config = ensure_config(config)

    def load_document(self, file_path: str | Path) -> list[Document]:
        """
        Load a document from file path.

        Args:
            file_path: Path to the document file

        Returns:
            List of loaded Document objects

        Raises:
            ValueError: If file type is not supported
            FileNotFoundError: If file doesn't exist
        """
        file_path = Path(file_path)

        if not file_path.exists():
            raise FileNotFoundError(f"File not found: {file_path}")

        file_ext = file_path.suffix.lower()

        if file_ext not in SUPPORTED_EXTENSIONS:
            raise ValueError(f"Unsupported file type: {file_ext}. Supported: {list(SUPPORTED_EXTENSIONS.keys())}")

        loader_class = DOCUMENT_LOADERS.get(file_ext)
        if not loader_class:
            raise ValueError(f"No loader available for {file_ext}")

        try:
            loader = loader_class(str(file_path))
            documents = loader.load()

            # Enrich metadata
            for doc in documents:
                if "file_name" not in doc.metadata:
                    doc.metadata["file_name"] = file_path.name
                if "file_type" not in doc.metadata:
                    doc.metadata["file_type"] = SUPPORTED_EXTENSIONS[file_ext]
                if "file_path" not in doc.metadata:
                    doc.metadata["file_path"] = str(file_path)
                if "source" not in doc.metadata:
                    doc.metadata["source"] = str(file_path)

            log_debug(f"Loaded {len(documents)} documents from {file_path}", self.config)
            return documents

        except Exception as e:
            log_error(f"Failed to load document {file_path}: {e}", self.config)
            raise

    def get_file_hash(self, file_path: str | Path) -> str:
        """Generate SHA-256 hash of file content"""
        hash_sha256 = hashlib.sha256()
        with open(file_path, "rb") as f:
            for chunk in iter(lambda: f.read(8192), b""):
                hash_sha256.update(chunk)
        return hash_sha256.hexdigest()

    def get_file_info(self, file_path: str | Path) -> DocumentMetadata:
        """Get comprehensive file information"""
        file_path = Path(file_path)

        return DocumentMetadata(
            file_path=str(file_path),
            file_name=file_path.name,
            file_type=SUPPORTED_EXTENSIONS.get(file_path.suffix.lower(), "Unknown"),
            file_size=file_path.stat().st_size,
            file_hash=self.get_file_hash(file_path),
            chunk_count=0,
            ingestion_timestamp=datetime.now().isoformat(),
            source_directory=str(file_path.parent),
        )


# ========== TEXT PREPROCESSING ==========
class TextPreprocessor:
    """Handles text normalization and deduplication"""

    def __init__(self, config: Optional["OrionConfig"] = None):
        self.config = ensure_config(config)
        self.preprocessing_config = self.config.rag.preprocessing
        self.similarity_threshold = self.preprocessing_config.similarity_threshold
        self.min_text_length = self.preprocessing_config.min_text_length

    def normalize_text(self, text: str) -> str:
        """
        Normalize text by cleaning and standardizing format.

        Args:
            text: Raw text to normalize

        Returns:
            Normalized text
        """
        if not text or not text.strip():
            return ""

        # Fix missing spaces
        text = self._fix_missing_spaces(text)

        # Remove excessive whitespace and normalize line endings
        text = re.sub(r"\s+", " ", text.strip())

        # Remove control characters but keep most printable characters
        text = re.sub(r"[\x00-\x08\x0B\x0C\x0E-\x1F\x7F]", "", text)

        # Normalize common text patterns
        text = re.sub(r"\.{2,}", "...", text)
        text = re.sub(r"\?{2,}", "?", text)
        text = re.sub(r"\!{2,}", "!", text)

        # Remove redundant spaces around punctuation
        text = re.sub(r"\s+([\.,:;!?])", r"\1", text)
        text = re.sub(r"([\.,:;!?])\s+", r"\1 ", text)

        # Normalize quotes
        text = re.sub(r'[""]', '"', text)
        text = re.sub(r"[''']", "'", text)

        return text.strip()

    def _fix_missing_spaces(self, text: str) -> str:
        """Fix missing spaces in PDF-extracted text using pattern recognition"""
        # Basic camelCase detection
        text = re.sub(r"([a-z])([A-Z])", r"\1 \2", text)

        # Numbers and letters
        text = re.sub(r"([a-zA-Z])(\d)", r"\1 \2", text)
        text = re.sub(r"(\d)([a-zA-Z])", r"\1 \2", text)

        # Common word boundaries
        text = re.sub(r"\bof([a-z]{3,})", r"of \1", text)
        text = re.sub(r"\bin([a-z]{3,})", r"in \1", text)
        text = re.sub(r"\bthe([a-z]{3,})", r"the \1", text)
        text = re.sub(r"\band([a-z]{3,})", r"and \1", text)

        # Clean up double spaces
        text = re.sub(r"\s+", " ", text)

        return text

    def calculate_similarity(self, text1: str, text2: str) -> float:
        """Calculate similarity between two texts"""
        if not text1 or not text2:
            return 0.0

        norm_text1 = self.normalize_text(text1).lower()
        norm_text2 = self.normalize_text(text2).lower()

        if not norm_text1 or not norm_text2:
            return 0.0

        matcher = SequenceMatcher(None, norm_text1, norm_text2)
        return matcher.ratio()

    def is_similar_content(self, text1: str, text2: str) -> bool:
        """Check if two texts are similar enough to be considered duplicates"""
        similarity = self.calculate_similarity(text1, text2)
        return similarity >= self.similarity_threshold

    def is_valid_text(self, text: str) -> bool:
        """Check if text meets minimum quality requirements"""
        if not text or not text.strip():
            return False

        if len(text.strip()) < self.min_text_length:
            return False

        meaningful_chars = re.sub(r"[^\w\s]", "", text)
        if len(meaningful_chars.strip()) < self.min_text_length * 0.7:
            return False

        return True

    def deduplicate_documents(self, documents: list[Document]) -> list[Document]:
        """Remove duplicate and similar documents"""
        if not documents or not self.preprocessing_config.enable_deduplication:
            return documents

        unique_documents = []
        seen_hashes = set()

        log_info(f"Deduplicating {len(documents)} documents (threshold: {self.similarity_threshold})", self.config)

        for doc in documents:
            # Quick hash-based deduplication
            content_hash = hashlib.sha256(doc.page_content.encode()).hexdigest()

            if content_hash in seen_hashes:
                log_debug("Skipped exact duplicate document", self.config)
                continue

            # Similarity-based deduplication
            is_duplicate = False
            for unique_doc in unique_documents:
                if self.is_similar_content(doc.page_content, unique_doc.page_content):
                    is_duplicate = True
                    log_debug("Skipped similar document (similarity check)", self.config)
                    break

            if not is_duplicate:
                unique_documents.append(doc)
                seen_hashes.add(content_hash)

        log_info(f"Deduplication: {len(documents)} -> {len(unique_documents)} documents", self.config)
        return unique_documents

    def preprocess_documents(self, documents: list[Document]) -> list[Document]:
        """Preprocess documents with normalization and deduplication"""
        if not documents:
            return documents

        # Normalize text if enabled
        if self.preprocessing_config.enable_normalization:
            for doc in documents:
                doc.page_content = self.normalize_text(doc.page_content)

        # Filter out invalid text
        valid_documents = [doc for doc in documents if self.is_valid_text(doc.page_content)]

        log_info(
            f"Text validation: {len(documents)} -> {len(valid_documents)} documents (min length: {self.min_text_length})",
            self.config,
        )

        # Deduplicate
        unique_documents = self.deduplicate_documents(valid_documents)

        return unique_documents


# ========== TEXT CHUNKING ==========
class TextChunker:
    """Handles different text chunking strategies"""

    def __init__(self, config: Optional["OrionConfig"] = None):
        self.config = ensure_config(config)
        self.chunking_config = self.config.rag.chunking

    def chunk_documents(self, documents: list[Document]) -> list[Document]:
        """
        Chunk documents using configured strategy.

        Args:
            documents: List of documents to chunk

        Returns:
            List of chunked documents
        """
        if not documents:
            return documents

        strategy = self.chunking_config.strategy

        log_info(
            f"Chunking {len(documents)} documents (strategy: {strategy.value}, "
            f"size: {self.chunking_config.chunk_size}, overlap: {self.chunking_config.chunk_overlap})",
            self.config,
        )

        if strategy == ChunkerType.RECURSIVE:
            return self._chunk_recursive(documents)
        elif strategy == ChunkerType.SEMANTIC:
            return self._chunk_semantic(documents)
        elif strategy == ChunkerType.SMART:
            return self._chunk_smart(documents)
        else:
            log_warning(f"Unknown chunking strategy: {strategy}, using recursive", self.config)
            return self._chunk_recursive(documents)

    def _chunk_recursive(self, documents: list[Document]) -> list[Document]:
        """Chunk using recursive character splitter"""
        splitter = RecursiveCharacterTextSplitter(
            chunk_size=self.chunking_config.chunk_size,
            chunk_overlap=self.chunking_config.chunk_overlap,
            length_function=len,
            separators=["\n\n", "\n", ". ", " ", ""],
        )

        chunked_docs = []
        for doc in documents:
            chunks = splitter.split_text(doc.page_content)

            for i, chunk in enumerate(chunks):
                # Clean chunk boundaries
                chunk = self._clean_chunk_boundaries(chunk)

                if len(chunk.strip()) < self.chunking_config.min_chunk_size:
                    continue

                # Create new document with chunk
                chunk_metadata = doc.metadata.copy()
                chunk_metadata["chunk_index"] = i
                chunk_metadata["chunk_count"] = len(chunks)

                chunked_docs.append(Document(page_content=chunk, metadata=chunk_metadata))

        log_info(f"Recursive chunking: {len(documents)} docs -> {len(chunked_docs)} chunks", self.config)
        return chunked_docs

    def _clean_chunk_boundaries(self, text: str) -> str:
        """Clean chunk boundaries to avoid mid-word splits"""
        text = text.strip()

        # Remove leading punctuation
        text = re.sub(r"^[^\w\s]+", "", text)

        # Remove trailing incomplete sentences
        if not text.endswith((".", "!", "?", ":", ";")):
            last_punct = max(
                text.rfind("."),
                text.rfind("!"),
                text.rfind("?"),
                text.rfind(":"),
                text.rfind(";"),
            )
            if last_punct > len(text) * 0.7:  # Keep if punctuation is in last 30%
                text = text[: last_punct + 1]

        return text.strip()

    def _chunk_semantic(self, documents: list[Document]) -> list[Document]:
        """Semantic chunking (placeholder for future implementation)"""
        log_warning("Semantic chunking not yet implemented, using recursive", self.config)
        return self._chunk_recursive(documents)

    def _chunk_smart(self, documents: list[Document]) -> list[Document]:
        """Smart chunking (placeholder for future implementation)"""
        log_warning("Smart chunking not yet implemented, using recursive", self.config)
        return self._chunk_recursive(documents)


# ========== MAIN INGESTION ORCHESTRATOR ==========
class DocumentIngestor:
    """Main orchestrator for document ingestion pipeline"""

    def __init__(
        self,
        vector_store: Optional[ChromaVectorStore] = None,
        config: Optional["OrionConfig"] = None,
    ):
        """
        Initialize document ingestor.

        Args:
            vector_store: Optional ChromaVectorStore instance. If None, creates new one.
            config: Orion configuration
        """
        self.config = ensure_config(config)
        self.processor = DocumentProcessor(config)
        self.preprocessor = TextPreprocessor(config)
        self.chunker = TextChunker(config)
        self.embedding_manager = EmbeddingManager(config=config)
        self.vector_store = vector_store or create_vector_store(config=config)

    def ingest_file(self, file_path: str | Path) -> tuple[bool, DocumentMetadata, list[str]]:
        """
        Ingest a single file into the vector store.

        Args:
            file_path: Path to file to ingest

        Returns:
            Tuple of (success, metadata, errors)
        """
        file_path = Path(file_path)
        errors = []

        try:
            # Get file info
            metadata = self.processor.get_file_info(file_path)

            # Load document
            documents = self.processor.load_document(file_path)

            if not documents:
                errors.append(f"No content extracted from {file_path}")
                return False, metadata, errors

            # Preprocess
            documents = self.preprocessor.preprocess_documents(documents)

            if not documents:
                errors.append(f"All content filtered out during preprocessing for {file_path}")
                return False, metadata, errors

            # Chunk
            chunks = self.chunker.chunk_documents(documents)

            if not chunks:
                errors.append(f"No chunks created for {file_path}")
                return False, metadata, errors

            # Extract text content from chunks
            chunk_texts = [chunk.page_content for chunk in chunks]

            # Generate embeddings
            log_debug(f"Generating embeddings for {len(chunk_texts)} chunks", self.config)
            embeddings = self.embedding_manager.encode_batch(chunk_texts)

            # Create metadata for each chunk
            chunk_metadatas = []
            for i, chunk in enumerate(chunks):
                chunk_metadata = self.vector_store.create_document_metadata(
                    file_path=file_path,
                    chunk_index=i,
                    total_chunks=len(chunks),
                    extra_metadata=chunk.metadata,  # Include any metadata from the Document object
                )
                chunk_metadatas.append(chunk_metadata)

            # Generate document IDs
            doc_ids = [f"{metadata.file_hash}_chunk_{i}" for i in range(len(chunks))]

            # Add to vector store
            success = self.vector_store.add_documents(
                documents=chunk_texts,
                embeddings=embeddings,
                metadatas=chunk_metadatas,
                ids=doc_ids,
            )

            if not success:
                errors.append("Failed to add documents to vector store")
                return False, metadata, errors

            # Update metadata
            metadata.chunk_count = len(chunks)
            metadata.ingestion_timestamp = datetime.now().isoformat()

            log_info(f"Successfully ingested {file_path}: {len(chunks)} chunks", self.config)
            return True, metadata, errors

        except Exception as e:
            error_msg = f"Failed to ingest {file_path}: {e}"
            log_error(error_msg, self.config)
            errors.append(error_msg)
            return False, metadata if "metadata" in locals() else None, errors

    def ingest_directory(
        self,
        directory_path: str | Path,
        recursive: bool = True,
        max_workers: int = 4,
    ) -> IngestionStats:
        """
        Ingest all supported files from a directory.

        Args:
            directory_path: Path to directory
            recursive: Whether to search recursively
            max_workers: Number of parallel workers

        Returns:
            IngestionStats with results
        """
        import time

        start_time = time.time()
        directory_path = Path(directory_path)

        if not directory_path.exists() or not directory_path.is_dir():
            log_error(f"Directory not found: {directory_path}", self.config)
            return IngestionStats()

        # Find all supported files
        files = []
        pattern = "**/*" if recursive else "*"

        for ext in SUPPORTED_EXTENSIONS.keys():
            files.extend(directory_path.glob(f"{pattern}{ext}"))

        if not files:
            log_warning(f"No supported files found in {directory_path}", self.config)
            return IngestionStats()

        log_info(f"Found {len(files)} files to ingest from {directory_path}", self.config)

        stats = IngestionStats(total_files=len(files))

        # Process files in parallel
        with ThreadPoolExecutor(max_workers=max_workers) as executor:
            futures = {executor.submit(self.ingest_file, file): file for file in files}

            with tqdm(total=len(files), desc="Ingesting files") as pbar:
                for future in as_completed(futures):
                    file = futures[future]
                    try:
                        success, metadata, errors = future.result()

                        if success:
                            stats.successful_files += 1
                            if metadata:
                                stats.total_chunks += metadata.chunk_count
                        else:
                            stats.failed_files += 1
                            for error in errors:
                                stats.add_error(error)

                    except Exception as e:
                        stats.failed_files += 1
                        error_msg = f"Unexpected error processing {file}: {e}"
                        stats.add_error(error_msg)
                        log_error(error_msg, self.config)

                    pbar.update(1)

        stats.processing_time = time.time() - start_time
        stats.total_embeddings = stats.total_chunks  # One embedding per chunk

        log_info(
            f"Ingestion complete: {stats.successful_files}/{stats.total_files} files, "
            f"{stats.total_chunks} chunks in {stats.processing_time:.2f}s",
            self.config,
        )

        return stats

    def ingest_knowledge_base(
        self,
        source_path: str | Path,
        clear_existing: bool = False,
    ) -> IngestionStats:
        """
        Ingest files from knowledge base path(s).

        Args:
            source_path: Path to file or directory
            clear_existing: Whether to clear existing data first

        Returns:
            IngestionStats with results
        """
        source_path = Path(source_path)

        if clear_existing:
            log_info("Clearing existing knowledge base", self.config)
            self.vector_store.clear_collection()

        if source_path.is_file():
            success, metadata, errors = self.ingest_file(source_path)
            stats = IngestionStats(total_files=1)
            if success:
                stats.successful_files = 1
                if metadata:
                    stats.total_chunks = metadata.chunk_count
            else:
                stats.failed_files = 1
                for error in errors:
                    stats.add_error(error)
            return stats

        elif source_path.is_dir():
            return self.ingest_directory(source_path, recursive=True)

        else:
            log_error(f"Path not found: {source_path}", self.config)
            return IngestionStats()

    def get_ingestion_summary(self) -> dict[str, Any]:
        """Get summary of current vector store state"""
        return self.vector_store.get_collection_stats()


# ========== CONVENIENCE FUNCTIONS ==========
def ingest_documents(
    source_path: str | Path,
    config: Optional["OrionConfig"] = None,
    clear_existing: bool = False,
) -> IngestionStats:
    """
    Convenience function to ingest documents.

    Args:
        source_path: Path to file or directory to ingest
        config: Optional Orion configuration
        clear_existing: Whether to clear existing data first

    Returns:
        IngestionStats with processing results

    Example:
        # Ingest a single file
        stats = ingest_documents("path/to/document.pdf")
        
        # Ingest entire directory
        stats = ingest_documents("path/to/knowledge_base")
        
        # Clear and re-ingest
        stats = ingest_documents("path/to/kb", clear_existing=True)
    """
    ingestor = DocumentIngestor(config=config)
    return ingestor.ingest_knowledge_base(source_path, clear_existing)


def get_supported_formats() -> dict[str, str]:
    """
    Get dictionary of supported file formats.

    Returns:
        Dictionary mapping file extensions to descriptions
    """
    return SUPPORTED_EXTENSIONS.copy()


def clear_knowledge_base(config: Optional["OrionConfig"] = None) -> bool:
    """
    Clear all data from the knowledge base.

    Args:
        config: Optional Orion configuration

    Returns:
        True if successful, False otherwise
    """
    try:
        vector_store = create_vector_store(config=config)
        return vector_store.clear_collection()
    except Exception as e:
        log_error(f"Failed to clear knowledge base: {e}", config)
        return False


def ingest_with_watchdog(
    watch_paths: Optional[list[str]] = None,
    config: Optional["OrionConfig"] = None,
) -> tuple[DocumentIngestor, Any]:
    """
    Create ingestor and file watcher for incremental ingestion.

    Args:
        watch_paths: List of paths to watch. If None, uses config.watchdog.paths
        config: Optional Orion configuration

    Returns:
        Tuple of (DocumentIngestor, FileWatcher)

    Example:
        # Start watching and auto-ingesting
        ingestor, watcher = ingest_with_watchdog(["/path/to/kb"])
        
        # Watcher automatically ingests new/modified files
        # Stop watching
        watcher.stop()
    """
    from src.retrieval.watchdog import create_file_watcher

    config = ensure_config(config)
    ingestor = DocumentIngestor(config=config)

    # Create callbacks for file events
    def handle_file_change(file_path: str):
        """Handle file addition or modification"""
        log_info(f"Ingesting file from watchdog: {file_path}", config)
        try:
            success, metadata, errors = ingestor.ingest_file(file_path)
            if success:
                log_info(f"Watchdog ingestion successful: {file_path}", config)
            else:
                log_error(f"Watchdog ingestion failed: {file_path} - {errors}", config)
        except Exception as e:
            log_error(f"Error during watchdog ingestion: {e}", config)

    # Create watcher with callbacks
    watcher = create_file_watcher(
        vector_store=ingestor.vector_store,
        on_file_added=handle_file_change,
        on_file_modified=handle_file_change,
        config=config,
    )

    # Start watching
    watcher.start(paths=watch_paths)

    return ingestor, watcher
